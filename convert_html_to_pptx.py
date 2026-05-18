import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme Definitions
COLOR_BG = RGBColor(28, 28, 33)        # #1C1C21 (Warm Dark Gray)
COLOR_WHITE = RGBColor(240, 240, 245)   # #F0F0F5 (Off-white)
COLOR_MUTED = RGBColor(160, 160, 170)   # #A0A0AA (Secondary text)
COLOR_GOLD = RGBColor(229, 192, 123)    # #E5C07B (Accent Gold)
COLOR_PEACH = RGBColor(255, 183, 178)   # #FFB7B2 (Accent Peach)
COLOR_BLUE = RGBColor(164, 187, 214)    # #A4BBD6 (Accent Blue)
COLOR_CARD_BG = RGBColor(38, 38, 43)    # Slightly lighter gray for card containers

# Helper: Set slide background color
def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

# Helper: Add unified top header block to a slide
def add_slide_header(slide, subtitle, title, right_desc=""):
    # Header container shape (thin border at bottom)
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.2))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Subtitle
    p_sub = tf.paragraphs[0]
    p_sub.text = subtitle.upper()
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(11)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_GOLD
    p_sub.space_after = Pt(4)
    
    # Main Title
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    
    if right_desc:
        right_box = slide.shapes.add_textbox(Inches(8.0), Inches(0.5), Inches(4.5), Inches(1.0))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        tf_right.margin_left = tf_right.margin_top = tf_right.margin_right = tf_right.margin_bottom = 0
        p_right = tf_right.paragraphs[0]
        p_right.text = right_desc
        p_right.alignment = PP_ALIGN.RIGHT
        p_right.font.name = "Arial"
        p_right.font.size = Pt(11)
        p_right.font.color.rgb = COLOR_MUTED

# Helper: Create standard info card background shape
def add_card_bg(slide, left, top, width, height, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# ==============================================================================
# SLIDE 1: Title Slide
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)

# Subtitle
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.0))
tf = title_box.text_frame
tf.word_wrap = True

p_sub = tf.paragraphs[0]
p_sub.alignment = PP_ALIGN.CENTER
p_sub.text = "PRECIOUS RECREATION"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(16)
p_sub.font.bold = True
p_sub.font.color.rgb = COLOR_GOLD
p_sub.space_after = Pt(20)

# Main Title
p_main = tf.add_paragraph()
p_main.alignment = PP_ALIGN.CENTER
p_main.text = "Life Cycle Memory Platform"
p_main.font.name = "Georgia"
p_main.font.size = Pt(48)
p_main.font.bold = True
p_main.font.color.rgb = COLOR_WHITE
p_main.space_after = Pt(14)

# Platform description
p_desc = tf.add_paragraph()
p_desc.alignment = PP_ALIGN.CENTER
p_desc.text = '"AI + 3D Hologram Technologies"'
p_desc.font.name = "Arial"
p_desc.font.size = Pt(20)
p_desc.font.bold = True
p_desc.font.color.rgb = COLOR_PEACH
p_desc.space_after = Pt(24)

# Footer tags
p_foot = tf.add_paragraph()
p_foot.alignment = PP_ALIGN.CENTER
p_foot.text = "Memory  •  Immersion  •  AI"
p_foot.font.name = "Arial"
p_foot.font.size = Pt(14)
p_foot.font.color.rgb = COLOR_BLUE

# ==============================================================================
# SLIDE 2: Founder & CEO Profile
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Founder & CEO Profile", "Sonny (YOUNGMIN SON)", "35 Years ICT / Mobile / Platform / Global Experiences")

# Profile Text Box
text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
tf = text_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

# Education
p1 = tf.paragraphs[0]
p1.text = "Education"
p1.font.name = "Arial"
p1.font.size = Pt(18)
p1.font.bold = True
p1.font.color.rgb = COLOR_GOLD
p1.space_after = Pt(6)

p1_sub = tf.add_paragraph()
p1_sub.text = "• Hanyang Univ (B.S. EE)\n• Univ. of Michigan (M.S. EE & CS)"
p1_sub.font.name = "Arial"
p1_sub.font.size = Pt(14)
p1_sub.font.color.rgb = COLOR_WHITE
p1_sub.space_after = Pt(20)

# Corporate Experience
p2 = tf.add_paragraph()
p2.text = "Professional Career"
p2.font.name = "Arial"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = COLOR_PEACH
p2.space_after = Pt(6)

p2_sub = tf.add_paragraph()
p2_sub.text = (
    "• SK Planet (2012 - 2026)\n"
    "  - PRM Leader, New Product Group Leader\n"
    "• SK Telecom (1999 - 2011)\n"
    "  - MMRC San Jose, SKT China Product Leader\n"
    "• Samsung Electronics (1993 - 1998)\n"
    "  - R&D Center Research Manager"
)
p2_sub.font.name = "Arial"
p2_sub.font.size = Pt(14)
p2_sub.font.color.rgb = COLOR_WHITE
p2_sub.space_after = Pt(14)

# Insert CEO Image (Right Column)
ceo_img_path = "ceo_profile_updated.png"
if os.path.exists(ceo_img_path):
    slide.shapes.add_picture(ceo_img_path, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8))
else:
    add_card_bg(slide, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8), COLOR_MUTED)

# ==============================================================================
# SLIDE 3: What is MBOX? (Floating Demo)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Platform Overview", "What is MBOX?", "AI + 3D 홀로그램 기반의 생애 주기 메모리 아카이빙 플랫폼")

# Left Column (overlapping images to simulate HTML composition)
# 1. 2D Photo Back Layer
wedding_img = "wedding_2d_input.jpg"
if os.path.exists(wedding_img):
    slide.shapes.add_picture(wedding_img, Inches(0.8), Inches(2.2), Inches(3.2), Inches(3.2))
else:
    add_card_bg(slide, Inches(0.8), Inches(2.2), Inches(3.2), Inches(3.2), COLOR_MUTED)

# Back layer label box
label_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(3.0), Inches(0.3))
p_lbl = label_box.text_frame.paragraphs[0]
p_lbl.text = "Flat 2D 원본 사진"
p_lbl.font.name = "Arial"
p_lbl.font.size = Pt(11)
p_lbl.font.color.rgb = COLOR_MUTED

# 2. Glowing Circle 3D Hologram Front Layer (The Spotlight center!)
holo_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.6), Inches(2.8), Inches(3.8), Inches(3.8))
holo_circle.fill.solid()
holo_circle.fill.fore_color.rgb = COLOR_BG
holo_circle.line.color.rgb = COLOR_GOLD
holo_circle.line.width = Pt(3.0)

# Add text or static image onto front layer
circle_tf = holo_circle.text_frame
circle_tf.word_wrap = True
p_c1 = circle_tf.paragraphs[0]
p_c1.alignment = PP_ALIGN.CENTER
p_c1.text = "\n\n[MBOX 3D HOLOGRAM]\n"
p_c1.font.name = "Arial"
p_c1.font.size = Pt(12)
p_c1.font.bold = True
p_c1.font.color.rgb = COLOR_GOLD

p_c2 = circle_tf.add_paragraph()
p_c2.alignment = PP_ALIGN.CENTER
p_c2.text = "3D Rotating Cube Video"
p_c2.font.name = "Arial"
p_c2.font.size = Pt(10)
p_c2.font.color.rgb = COLOR_WHITE

# Right Column Text Box
desc_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
tf = desc_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

p_rh = tf.paragraphs[0]
p_rh.text = "AI-powered automated video generation platform"
p_rh.font.name = "Arial"
p_rh.font.size = Pt(18)
p_rh.font.bold = True
p_rh.font.color.rgb = COLOR_PEACH
p_rh.space_after = Pt(10)

p_rd = tf.add_paragraph()
p_rd.text = "스마트폰 속 잠들어 있는 평범한 사진들을 공간적 입체감이 살아 숨쉬는 고품질 3D 홀로그램 비디오 에셋으로 단 몇 분 만에 자동 변환합니다."
p_rd.font.name = "Arial"
p_rd.font.size = Pt(13)
p_rd.font.color.rgb = COLOR_WHITE
p_rd.space_after = Pt(20)

# Capabilities Box
add_card_bg(slide, Inches(6.8), Inches(3.6), Inches(5.7), Inches(3.0), COLOR_GOLD)
cap_box = slide.shapes.add_textbox(Inches(7.1), Inches(3.8), Inches(5.1), Inches(2.6))
tf_cap = cap_box.text_frame
tf_cap.word_wrap = True
tf_cap.margin_left = tf_cap.margin_top = tf_cap.margin_right = tf_cap.margin_bottom = 0

p_cap_t = tf_cap.paragraphs[0]
p_cap_t.text = "핵심 역량 (Core Capabilities)"
p_cap_t.font.name = "Arial"
p_cap_t.font.size = Pt(15)
p_cap_t.font.bold = True
p_cap_t.font.color.rgb = COLOR_GOLD
p_cap_t.space_after = Pt(8)

p_cap_l = tf_cap.add_paragraph()
p_cap_l.text = (
    "• Auto Format: 빠른 콘텐츠 어셈블리를 위한 해상도 자동 규격화\n\n"
    "• 3D Synthesis: 딥러닝 기반 배경 분리 및 3D 공간 레이어 합성\n\n"
    "• Hologram Delivery: 전용 H30 홀로그램 디바이스 즉시 송출 통합"
)
p_cap_l.font.name = "Arial"
p_cap_l.font.size = Pt(12)
p_cap_l.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 4: Business Model (3 Segmentation)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Business Model", "3 Market Segmentation", "B2B, B2B2C, B2C를 아우르는 다각적 플랫폼 비즈니스 모델")

col_w = Inches(3.6)
col_h = Inches(4.6)
top_y = Inches(1.8)

# Col 1: B2B
add_card_bg(slide, Inches(0.8), top_y, col_w, col_h, COLOR_BLUE)
c1_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(4.2))
tf1 = c1_box.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
p = tf1.paragraphs[0]
p.text = "B2B MARKET"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(4)

p_sub = tf1.add_paragraph()
p_sub.text = "Enterprise & Government"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(12)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(20)

p_body = tf1.add_paragraph()
p_body.text = (
    "• 타겟: 박물관, 전시관, 기업 홍보관 등\n\n"
    "• 내용: 3D 미디어 콘텐츠 제작 및 H30 홀로그램 딜리버리 시스템 통합\n\n"
    "• 수익: 프로젝트 단위 라이선스 및 유지관리 하드웨어 수수료"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Col 2: B2B2C
add_card_bg(slide, Inches(4.8), top_y, col_w, col_h, COLOR_GOLD)
c2_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.2))
tf2 = c2_box.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
p = tf2.paragraphs[0]
p.text = "B2B2C MARKET"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(4)

p_sub = tf2.add_paragraph()
p_sub.text = "Strategic Partner Integration"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(12)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(20)

p_body = tf2.add_paragraph()
p_body.text = (
    "• 타겟: 반려동물 추모공원, 웨딩홀 등\n\n"
    "• 내용: 로비 공간 내 홀로그램 키오스크 서비스 무인 운영 연계\n\n"
    "• 수익: 고객 유료 서비스 구매 시 파트너사와 50:50 건당 분배"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Col 3: B2C
add_card_bg(slide, Inches(8.8), top_y, col_w, col_h, COLOR_PEACH)
c3_box = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(4.2))
tf3 = c3_box.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
p = tf3.paragraphs[0]
p.text = "B2C MARKET"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(4)

p_sub = tf3.add_paragraph()
p_sub.text = "Direct-to-Consumer App"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(12)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(20)

p_body = tf3.add_paragraph()
p_body.text = (
    "• 타겟: 개인 아카이빙 회원\n\n"
    "• 내용: MBOX 전용 모바일 앱을 통한 개인 소장용 홀로그램 콘텐츠 생성\n\n"
    "• 수익: 클라우드 보관 수수료 및 고품질 에셋 렌더링 구독 월회비"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 5: TAM SAM SOM Market Size
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Market Opportunity", "Life Cycle Memory B2B2C TAM", "반려동물 추모 및 웨딩 플랫폼 타겟 시장 규모 산출")

# Left details
market_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
tf = market_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

# TAM
p_tam = tf.paragraphs[0]
p_tam.text = "B2C TAM : 4,325억 원"
p_tam.font.name = "Arial"
p_tam.font.size = Pt(18)
p_tam.font.bold = True
p_tam.font.color.rgb = COLOR_PEACH
p_tam.space_after = Pt(4)

p_tam_d = tf.add_paragraph()
p_tam_d.text = "대한민국 인구 5천만 명 × 생애 전환기 서비스 이용 연간 LTV 60만 원 타겟 시장 규모"
p_tam_d.font.name = "Arial"
p_tam_d.font.size = Pt(12.5)
p_tam_d.font.color.rgb = COLOR_WHITE
p_tam_d.space_after = Pt(16)

# SAM
p_sam = tf.add_paragraph()
p_sam.text = "B2B2C SAM : 2,100억 원"
p_sam.font.name = "Arial"
p_sam.font.size = Pt(18)
p_sam.font.bold = True
p_sam.font.color.rgb = COLOR_GOLD
p_sam.space_after = Pt(4)

p_sam_d = tf.add_paragraph()
p_sam_d.text = "전국 반려동물 추모공원 12.1만 안치, 장례 35만, 연간 예식 24만 건 기반 결합 디지털 사이니지 시장 규모"
p_sam_d.font.name = "Arial"
p_sam_d.font.size = Pt(12.5)
p_sam_d.font.color.rgb = COLOR_WHITE
p_sam_d.space_after = Pt(16)

# SOM
p_som = tf.add_paragraph()
p_som.text = "B2B SOM : 5,075억 원"
p_som.font.name = "Arial"
p_som.font.size = Pt(18)
p_som.font.bold = True
p_som.font.color.rgb = COLOR_BLUE
p_som.space_after = Pt(4)

p_som_d = tf.add_paragraph()
p_som_d.text = "국내 디지털 미디어 사이니지 시장 규모 1.3조 원 중 플랫폼 딜리버리 및 장비 락인 실현 가능 초기 점유 시장"
p_som_d.font.name = "Arial"
p_som_d.font.size = Pt(12.5)
p_som_d.font.color.rgb = COLOR_WHITE

# Right Segment Donut Mockup
donut_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), Inches(2.3), Inches(3.8), Inches(3.8))
donut_bg.fill.solid()
donut_bg.fill.fore_color.rgb = COLOR_CARD_BG
donut_bg.line.color.rgb = COLOR_PEACH
donut_bg.line.width = Pt(4.0)

# Inner hole to make it look like a donut
donut_hole = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(3.2), Inches(2.0), Inches(2.0))
donut_hole.fill.solid()
donut_hole.fill.fore_color.rgb = COLOR_BG
donut_hole.line.fill.background()

# Title in center of donut
inner_tf = donut_hole.text_frame
inner_tf.word_wrap = True
p_in = inner_tf.paragraphs[0]
p_in.alignment = PP_ALIGN.CENTER
p_in.text = "\nPlatform TAM\n1.15조"
p_in.font.name = "Arial"
p_in.font.size = Pt(13)
p_in.font.bold = True
p_in.font.color.rgb = COLOR_GOLD

# ==============================================================================
# SLIDE 6: AI Video Technology Pipeline
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "AI Engine Pipeline", "AI Video Technology Pipeline", "스마트폰의 평면 사진을 3D 홀로그램 비디오 에셋으로 자동화 연산하는 핵심 파이프라인")

steps = [
    ("1", "Photo Upload", "스마트폰 2D 원본 사진 업로드"),
    ("2", "AI Segment", "피사체 추출 및 배경 실시간 분리"),
    ("3", "Template", "H30 디바이스 전용 1024×1024 규격화"),
    ("4", "BG Synthesis", "3D 깊이 맵(Depth Map) 자동 연산 및 합성"),
    ("5", "Rendering", "30fps 실시간 고품질 비디오 인코딩"),
    ("6", "Distribution", "H30 디바이스 즉시 송출 배포")
]

card_w = Inches(1.6)
card_h = Inches(3.2)
start_x = Inches(0.8)
gap_x = Inches(2.0)
y_pos = Inches(2.0)

for i, (num, title, body) in enumerate(steps):
    x = start_x + (i * gap_x)
    # Background shape
    add_card_bg(slide, x, y_pos, card_w, card_h, COLOR_PEACH)
    
    # Step bubble
    bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), y_pos + Inches(0.3), Inches(0.6), Inches(0.6))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = COLOR_PEACH
    bubble.line.fill.background()
    p_b = bubble.text_frame.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    p_b.text = num
    p_b.font.name = "Arial"
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_BG
    
    # Content Text Box
    step_box = slide.shapes.add_textbox(x + Inches(0.1), y_pos + Inches(1.1), card_w - Inches(0.2), card_h - Inches(1.3))
    tf = step_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_t = tf.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    p_t.text = title
    p_t.font.name = "Arial"
    p_t.font.size = Pt(12.5)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_WHITE
    p_t.space_after = Pt(6)
    
    p_d = tf.add_paragraph()
    p_d.alignment = PP_ALIGN.CENTER
    p_d.text = body
    p_d.font.name = "Arial"
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = COLOR_MUTED
    
    # Arrow (draw for steps 1-5)
    if i < 5:
        arrow_box = slide.shapes.add_textbox(x + card_w, y_pos + Inches(1.3), Inches(0.4), Inches(0.6))
        p_arr = arrow_box.text_frame.paragraphs[0]
        p_arr.alignment = PP_ALIGN.CENTER
        p_arr.text = "➡️"
        p_arr.font.name = "Arial"
        p_arr.font.size = Pt(14)
        p_arr.font.color.rgb = COLOR_GOLD

# Tech warning note at bottom
add_card_bg(slide, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), COLOR_PEACH)
note_box = slide.shapes.add_textbox(Inches(1.1), Inches(5.75), Inches(11.133), Inches(0.9))
tf_note = note_box.text_frame
tf_note.word_wrap = True
tf_note.margin_left = tf_note.margin_top = tf_note.margin_right = tf_note.margin_bottom = 0

p_n = tf_note.paragraphs[0]
p_n.text = "💡 로컬 우선(Local-First) 자동화 연산:"
p_n.font.name = "Arial"
p_n.font.size = Pt(12)
p_n.font.bold = True
p_n.font.color.rgb = COLOR_PEACH
p_n.space_after = Pt(2)

p_n_d = tf_note.add_paragraph()
p_n_d.text = "MBOX AI 엔진 파이프라인은 모든 연산 과정을 네트워크 레이텐시 없이 디바이스 로컬에서 수 분 내에 처리하도록 최적화되어 있어, 상용화 즉시 비용 효율적인 대규모 확장이 가능합니다."
p_n_d.font.name = "Arial"
p_n_d.font.size = Pt(11)
p_n_d.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 7: Life's Every Moment Category
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Market Vertical", "Life's Every Moment : 생애 주기 카테고리", "인간 생애 전반의 소중한 로맨틱 마일스톤과 추모 패밀리 아카이브")

col_w = Inches(3.6)
col_h = Inches(4.5)
y_pos = Inches(1.8)

# Cat 1: Memorial
add_card_bg(slide, Inches(0.8), y_pos, col_w, col_h, COLOR_PEACH)
cat1_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf1 = cat1_box.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

p = tf1.paragraphs[0]
p.text = "Memorial Album"
p.font.name = "Arial"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(8)

p_sub = tf1.add_paragraph()
p_sub.text = "반려동물 및 장례 추모 헌정"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(13)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(14)

p_body = tf1.add_paragraph()
p_body.text = (
    "• 반려견/반려묘와의 추억 공간 복원\n\n"
    "• 장례식장 및 납골당 내 입체 추모관 연동\n\n"
    "• 이별의 슬픔을 넘어선 따뜻한 디지털 헌정 앨범 제공"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Cat 2: Love & Wedding
add_card_bg(slide, Inches(4.8), y_pos, col_w, col_h, COLOR_GOLD)
cat2_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf2 = cat2_box.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0

p = tf2.paragraphs[0]
p.text = "Love & Wedding"
p.font.name = "Arial"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(8)

p_sub = tf2.add_paragraph()
p_sub.text = "로맨틱 마일스톤 및 웨딩"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(13)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(14)

p_body = tf2.add_paragraph()
p_body.text = (
    "• 결혼식장 로비 3D 웰컴 보드 시청\n\n"
    "• 신랑신부 프로포즈 및 러브스토리 입체 복원\n\n"
    "• 하객들에게 생생하고 차별화된 시각적 공간 연출 제공"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Cat 3: Family & Travel
add_card_bg(slide, Inches(8.8), y_pos, col_w, col_h, COLOR_BLUE)
cat3_box = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf3 = cat3_box.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0

p = tf3.paragraphs[0]
p.text = "Family & Travel"
p.font.name = "Arial"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(8)

p_sub = tf3.add_paragraph()
p_sub.text = "가족 연대기 및 여행 저널"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(13)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(14)

p_body = tf3.add_paragraph()
p_body.text = (
    "• 부모님 칠순/팔순 기념 아카이브\n\n"
    "• 아이 성장 일기 및 가족 여행 입체 렌더링\n\n"
    "• 평범한 거실 공간을 몰입형 생애 미술관으로 확장"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 8: Pet Use Case (Dog)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Use Case 01", "반려동물 추모 서비스 : 강아지 (Pet Memorial)", "추모공원 H30 디바이스를 통한 감성 펫로스 치유 및 아카이브 연동")

# Left Column text
details_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
tf = details_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

# Challenge
p = tf.paragraphs[0]
p.text = "THE CHALLENGE (시장 통증)"
p.font.name = "Arial"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(2)

p_d = tf.add_paragraph()
p_d.text = "국내 반려인 1,500만 시대, 반려견의 죽음으로 인한 '펫로스 증후군(Pet Loss Syndrome)'은 심각한 사회적 현상이지만 감성적인 추모 솔루션이 부재합니다."
p_d.font.name = "Arial"
p_d.font.size = Pt(12)
p_d.font.color.rgb = COLOR_WHITE
p_d.space_after = Pt(12)

# Solution
p = tf.add_paragraph()
p.text = "OUR SOLUTION (MBOX 해결책)"
p.font.name = "Arial"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(2)

p_d = tf.add_paragraph()
p_d.text = "추모공원 내 H30 입체 홀로그램 디스플레이 키오스크를 도입, AI 엔진이 반려견 평면 사진을 분석하여 12가지 생생한 감정선 맞춤 입체 콘텐츠를 복원 상영합니다."
p_d.font.name = "Arial"
p_d.font.size = Pt(12)
p_d.font.color.rgb = COLOR_WHITE
p_d.space_after = Pt(12)

# Differentiation
p = tf.add_paragraph()
p.text = "DIFFERENTIATION (차별화 요소)"
p.font.name = "Arial"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(2)

p_d = tf.add_paragraph()
p_d.text = "국내 최초의 3D 홀로그램 펫 추모 기술력, 반려동물 납골당과 연계하는 B2B2C 파트너십 상생 비즈니스 모델, 24시간 무인 키오스크 자동화 시스템 및 감성 반응형 AI를 실현합니다."
p_d.font.name = "Arial"
p_d.font.size = Pt(12)
p_d.font.color.rgb = COLOR_WHITE

# Right Column Image
pet_img = "dog_hologram_memorial.png"
if os.path.exists(pet_img):
    slide.shapes.add_picture(pet_img, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8))
else:
    add_card_bg(slide, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8), COLOR_PEACH)

# ==============================================================================
# SLIDE 9: Pet Revenue Projections
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Go-To-Market", "Pet Memorial B2B2C 수익 시뮬레이션", "전국 반려동물 추모공원(납골당) 파트너십 구축 및 유료 연회원(Renewal) 기반 성장 모델")

# Projections columns
col_w = Inches(3.6)
col_h = Inches(3.2)
top_y = Inches(1.8)

# 2026
add_card_bg(slide, Inches(0.8), top_y, col_w, col_h, COLOR_GOLD)
p_box1 = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(2.8))
tf1 = p_box1.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
p = tf1.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "1.57억 원"
p.font.name = "Arial"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(14)
p_body = tf1.add_paragraph()
p_body.alignment = PP_ALIGN.CENTER
p_body.text = "2026 (도입기)\n\n• 파트너 납골당 5곳 확보\n• 유료 아카이빙 회원 350명 달성"
p_body.font.name = "Arial"
p_body.font.size = Pt(13)
p_body.font.color.rgb = COLOR_WHITE

# 2027
add_card_bg(slide, Inches(4.8), top_y, col_w, col_h, COLOR_PEACH)
p_box2 = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(2.8))
tf2 = p_box2.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "6.56억 원"
p.font.name = "Arial"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(14)
p_body = tf2.add_paragraph()
p_body.alignment = PP_ALIGN.CENTER
p_body.text = "2027 (성장기)\n\n• 파트너 납골당 20곳 확장\n• 갱신 포함 회원 1,575명 누적"
p_body.font.name = "Arial"
p_body.font.size = Pt(13)
p_body.font.color.rgb = COLOR_WHITE

# 2028
add_card_bg(slide, Inches(8.8), top_y, col_w, col_h, COLOR_BLUE)
p_box3 = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(2.8))
tf3 = p_box3.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "16.9억 원"
p.font.name = "Arial"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(14)
p_body = tf3.add_paragraph()
p_body.alignment = PP_ALIGN.CENTER
p_body.text = "2028 (확장기)\n\n• 파트너 납골당 50곳 선점\n• 누적 갱신 회원 4,287명 락인"
p_body.font.name = "Arial"
p_body.font.size = Pt(13)
p_body.font.color.rgb = COLOR_WHITE

# Bottom Details
add_card_bg(slide, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.6))
details_box = slide.shapes.add_textbox(Inches(1.1), Inches(5.4), Inches(11.133), Inches(1.2))
tf_det = details_box.text_frame
tf_det.word_wrap = True
tf_det.margin_left = tf_det.margin_top = tf_det.margin_right = tf_det.margin_bottom = 0

# 3 Columns inside bottom card
col1 = tf_det.paragraphs[0]
col1.text = "📊 산출 근거 (Basis): 반려동물 추모공원 납골당 평균 안치 1,600건 기준, 잠재 고객 30%(500명) 타겟팅"
col1.font.name = "Arial"
col1.font.size = Pt(11.5)
col1.font.color.rgb = COLOR_GOLD
col1.space_after = Pt(6)

col2 = tf_det.add_paragraph()
col2.text = "🔄 유료 전환율 (Conversion): 체험 고객 중 35% 연회원 전환 가정 (Basic 30만 원 / Premium 50만 원 구성)"
col2.font.name = "Arial"
col2.font.size = Pt(11.5)
col2.font.color.rgb = COLOR_PEACH
col2.space_after = Pt(6)

col3 = tf_det.add_paragraph()
col3.text = "📈 갱신 및 재구독 (Renewal): 누적 갱신 회원의 지속적 추가로 시간이 갈수록 기하급수적으로 이익률 상승"
col3.font.name = "Arial"
col3.font.size = Pt(11.5)
col3.font.color.rgb = COLOR_BLUE

# ==============================================================================
# SLIDE 10: Wedding Use Case
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Use Case 02", "웨딩 메모리얼 서비스 (Wedding Memorial)", "웨딩홀 로비 MBOX 홀로그램 보드로 하객 대기 공간 경험 및 가치 혁신")

# Left Column Image
wedding_holo_img = "wedding_hologram_memorial.png"
if os.path.exists(wedding_holo_img):
    slide.shapes.add_picture(wedding_holo_img, Inches(0.8), Inches(1.8), Inches(4.7), Inches(4.8))
else:
    add_card_bg(slide, Inches(0.8), Inches(1.8), Inches(4.7), Inches(4.8), COLOR_GOLD)

# Right Column text
details_box = slide.shapes.add_textbox(Inches(6.0), Inches(1.8), Inches(6.5), Inches(5.0))
tf = details_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

# Concept
p = tf.paragraphs[0]
p.text = "CONCEPT (공간 경험의 혁신)"
p.font.name = "Arial"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(2)

p_d = tf.add_paragraph()
p_d.text = "웨딩홀 로비에서 하객들이 대기하는 지루한 약 30분의 시간을 H30 입체 홀로그램 디스플레이를 활용한 '두 사람의 가장 아름답고 눈부신 러브스토리 시청의 순간'으로 재창조하여 품격을 높입니다."
p_d.font.name = "Arial"
p_d.font.size = Pt(12)
p_d.font.color.rgb = COLOR_WHITE
p_d.space_after = Pt(12)

# Content Creation
p = tf.add_paragraph()
p.text = "CONTENT CREATION (완전 자동화)"
p.font.name = "Arial"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(2)

p_d = tf.add_paragraph()
p_d.text = "신랑신부가 제출한 20~40장의 일반 평면 스마트폰 사진을 MBOX AI 엔진이 1024×1024 포맷의 3D 입체 에셋으로 변환하여, 각 사진당 3~5초의 트랜지션과 입체 배경을 합성한 30분 분량의 로비 대기 특화 재생 목록을 즉시 생성합니다."
p_d.font.name = "Arial"
p_d.font.size = Pt(12)
p_d.font.color.rgb = COLOR_WHITE
p_d.space_after = Pt(12)

# Business Value
p = tf.add_paragraph()
p.text = "BUSINESS VALUE (비즈니스적 윈윈)"
p.font.name = "Arial"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(2)

p_d = tf.add_paragraph()
p_d.text = "웨딩홀 파트너는 추가 인력이나 별도 편집 수고 없이 무인으로 시스템을 운영하며, MBOX와의 50:50 수익 배분을 통해 연간 약 5천만 원 이상의 개별 고부가가치 부가 수익을 창출하는 완벽한 윈윈 비즈니스를 실현합니다."
p_d.font.name = "Arial"
p_d.font.size = Pt(12)
p_d.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 11: Wedding Revenue Projections
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Go-To-Market", "Wedding Memorial B2B2C 수익 시뮬레이션", "전국 웨딩홀 파트너십 구축 및 예식당 건별 과금(Pay-per-use) + 50:50 수익 배분 성장 모델")

# Projections columns (using matching layout but wedding colors/math)
# 2026
add_card_bg(slide, Inches(0.8), top_y, col_w, col_h, COLOR_GOLD)
p_box1 = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(2.8))
tf1 = p_box1.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
p = tf1.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "1.8억 원"
p.font.name = "Arial"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(14)
p_body = tf1.add_paragraph()
p_body.alignment = PP_ALIGN.CENTER
p_body.text = "2026 (도입기)\n\n• 파트너 웨딩홀 20곳 확보\n• 연간 1,200건 예식 적용\n• (MBOX Net Share: 50% 분배)"
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# 2027
add_card_bg(slide, Inches(4.8), top_y, col_w, col_h, COLOR_PEACH)
p_box2 = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(2.8))
tf2 = p_box2.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "7.2억 원"
p.font.name = "Arial"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(14)
p_body = tf2.add_paragraph()
p_body.alignment = PP_ALIGN.CENTER
p_body.text = "2027 (성장기)\n\n• 파트너 웨딩홀 60곳 확장\n• 연간 4,800건 예식 적용\n• (무인 자동화 정산 모델 확장)"
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# 2028
add_card_bg(slide, Inches(8.8), top_y, col_w, col_h, COLOR_BLUE)
p_box3 = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(2.8))
tf3 = p_box3.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "18.0억 원"
p.font.name = "Arial"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(14)
p_body = tf3.add_paragraph()
p_body.alignment = PP_ALIGN.CENTER
p_body.text = "2028 (확장기)\n\n• 파트너 웨딩홀 150곳 선점\n• 연간 12,000건 달성 (M/S 5%)\n• 연 18억 독점적 플랫폼 매출"
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Bottom Details
add_card_bg(slide, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.6))
details_box = slide.shapes.add_textbox(Inches(1.1), Inches(5.4), Inches(11.133), Inches(1.2))
tf_det = details_box.text_frame
tf_det.word_wrap = True
tf_det.margin_left = tf_det.margin_top = tf_det.margin_right = tf_det.margin_bottom = 0

# 3 Columns inside bottom card
col1 = tf_det.paragraphs[0]
col1.text = "📊 산출 근거 (Basis): 전국 웨딩홀 평균 예식 수 연 320건 기준, 대기 하객 대상 MBOX 패키지 30%(96건) 도입 타겟"
col1.font.name = "Arial"
col1.font.size = Pt(11.5)
col1.font.color.rgb = COLOR_GOLD
col1.space_after = Pt(6)

col2 = tf_det.add_paragraph()
col2.text = "🔄 수익 배분 모델 (Split): 예식당 30만 원 패키지 적용 시, 웨딩홀 파트너와 MBOX 간의 50:50 수익 정산 모델 수립"
col2.font.name = "Arial"
col2.font.size = Pt(11.5)
col2.font.color.rgb = COLOR_PEACH
col2.space_after = Pt(6)

col3 = tf_det.add_paragraph()
col3.text = "📈 확장 시너지 (Synergy): 전국 733개 웨딩홀 network로 신속한 락인(Lock-in) 효과 발생 및 플랫폼 기반 웨딩 생태계 확장"
col3.font.name = "Arial"
col3.font.size = Pt(11.5)
col3.font.color.rgb = COLOR_BLUE

# ==============================================================================
# SLIDE 12: Core Technology
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Technical Foundation", "Core Technology", "MBOX 플랫폼의 지속적인 경쟁력을 유지하는 핵심 기술 요소")

# Three columns
# Tech 1
add_card_bg(slide, Inches(0.8), y_pos, col_w, col_h, COLOR_GOLD)
t1_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf1 = t1_box.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
p = tf1.paragraphs[0]
p.text = "3D Hologram H30"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(14)
p_body = tf1.add_paragraph()
p_body.text = (
    "• 1024×1024 고해상도\n\n"
    "• 24fps H.264 High + AAC 오디오 지원\n\n"
    "• 1:1 몰입형 디스플레이 레이아웃 설계\n\n"
    "• HoloVision H30 전용 하드웨어 연동"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(13)
p_body.font.color.rgb = COLOR_WHITE

# Tech 2
add_card_bg(slide, Inches(4.8), y_pos, col_w, col_h, COLOR_BLUE)
t2_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf2 = t2_box.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
p = tf2.paragraphs[0]
p.text = "Local-First AI"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(14)
p_body = tf2.add_paragraph()
p_body.text = (
    "• Whisper.cpp (STT) 음성 인식 최적화\n\n"
    "• Llama.cpp / Ollama 로컬 LLM 통합\n\n"
    "• On-device 처리를 통한 고객 데이터 누출 원천 차단\n\n"
    "• 외부망 장애 시에도 무중단 안정적 상영"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(13)
p_body.font.color.rgb = COLOR_WHITE

# Tech 3
add_card_bg(slide, Inches(8.8), y_pos, col_w, col_h, COLOR_PEACH)
t3_box = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf3 = t3_box.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
p = tf3.paragraphs[0]
p.text = "12-Emotion Routing"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(14)
p_body = tf3.add_paragraph()
p_body.text = (
    "• 12가지의 감정 스펙트럼 세밀 분석\n\n"
    "• 감정 매칭 알고리즘을 통한 맞춤형 테마 및 음악 매칭\n\n"
    "• 슬픔, 따뜻함, 행복 등 하객 반응형 AI 인터랙션 제공"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(13)
p_body.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 13: Core Engineering Expertise
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Engineering Organization", "Core Engineering Expertise", "상용화 수준의 딥테크를 최적 설계하는 MBOX 핵심 엔지니어 역량")

# Three columns
# Eng 1
add_card_bg(slide, Inches(0.8), y_pos, col_w, col_h, COLOR_BLUE)
e1_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf1 = e1_box.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
p = tf1.paragraphs[0]
p.text = "AI Vision Lab"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE
p.space_after = Pt(4)
p_sub = tf1.add_paragraph()
p_sub.text = "Generative AI & LLM Experts"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(11)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(14)
p_body = tf1.add_paragraph()
p_body.text = (
    "• 단순 API 연동이 아닌 로컬 모델 자체 경량화/최적화\n\n"
    "• 1024×1024 해상도 피사체 실시간 분리 및 깊이 맵 합성\n\n"
    "• 12-감정선 분석 라우팅 독자 알고리즘 설계"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Eng 2
add_card_bg(slide, Inches(4.8), y_pos, col_w, col_h, COLOR_GOLD)
e2_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf2 = e2_box.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
p = tf2.paragraphs[0]
p.text = "Media Engine Team"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(4)
p_sub = tf2.add_paragraph()
p_sub.text = "3D Rendering & Video Processing"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(11)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(14)
p_body = tf2.add_paragraph()
p_body.text = (
    "• Zero-latency 비디오 자동 생성 파이프라인 최적 설계\n\n"
    "• 24fps H.264 High 인코딩 무인 자동화 서버 구축\n\n"
    "• 대용량 홀로그램 멀티채널 배포 인프라 및 CDN 아키텍처"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# Eng 3
add_card_bg(slide, Inches(8.8), y_pos, col_w, col_h, COLOR_PEACH)
e3_box = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(4.1))
tf3 = e3_box.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
p = tf3.paragraphs[0]
p.text = "Hardware Infra"
p.font.name = "Arial"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PEACH
p.space_after = Pt(4)
p_sub = tf3.add_paragraph()
p_sub.text = "IoT & Embedded Systems"
p_sub.font.name = "Arial"
p_sub.font.size = Pt(11)
p_sub.font.color.rgb = COLOR_MUTED
p_sub.space_after = Pt(14)
p_body = tf3.add_paragraph()
p_body.text = (
    "• HoloVision H30 전용 광학/드라이버 하드웨어 직접 설계\n\n"
    "• 24시간 무인 운영을 위한 온도/무결점 제어 회로 설계\n\n"
    "• 원격 모니터링 및 OTA 보안 펌웨어 라이브러리"
)
p_body.font.name = "Arial"
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = COLOR_WHITE

# ==============================================================================
# SLIDE 14: BI Center Cooperation & Synergy
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
add_slide_header(slide, "Cooperation & Ecosystem", "BI Center × MBOX × HoloVision Synergy", "산학 협력을 통한 R&D 기술 고도화 및 미래 AI 서비스 확장 생태계")

# Left Column (Synergy Diagram represented nicely in text)
diag_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.0), Inches(5.0))
tf_diag = diag_box.text_frame
tf_diag.word_wrap = True
tf_diag.margin_left = tf_diag.margin_top = tf_diag.margin_right = tf_diag.margin_bottom = 0

p_dg = tf_diag.paragraphs[0]
p_dg.text = "MBOX SYNERGY CIRCLE"
p_dg.font.name = "Arial"
p_dg.font.size = Pt(18)
p_dg.font.bold = True
p_dg.font.color.rgb = COLOR_GOLD
p_dg.space_after = Pt(20)

p_dg_b = tf_diag.add_paragraph()
p_dg_b.text = (
    "🌐 BI Center (Hanshin Univ.)\n"
    "  - 인큐베이션 공간 제공, 특허 매핑, 국책 R&D 프로젝트 발굴\n\n"
    "🤖 MBOX (Core AI Engine)\n"
    "  - AI Vision, 3D 합성, 12-감정선 라우팅 알고리즘 개발 주도\n\n"
    "🖥️ HoloVision (Display HW)\n"
    "  - H30 전용 광학 디바이스 및 하드웨어 모니터링 인프라 양산"
)
p_dg_b.font.name = "Arial"
p_dg_b.font.size = Pt(14)
p_dg_b.font.color.rgb = COLOR_WHITE

# Right Column AS-IS & TO-BE Cards
# AS-IS Card
add_card_bg(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.2), COLOR_PEACH)
asis_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.3), Inches(1.9))
tf_as = asis_box.text_frame
tf_as.word_wrap = True
tf_as.margin_left = tf_as.margin_top = tf_as.margin_right = tf_as.margin_bottom = 0

p_as_t = tf_as.paragraphs[0]
p_as_t.text = "AS-IS: 현 협력 현황 (Current Progress)"
p_as_t.font.name = "Arial"
p_as_t.font.size = Pt(14)
p_as_t.font.bold = True
p_as_t.font.color.rgb = COLOR_PEACH
p_as_t.space_after = Pt(6)

p_as_d = tf_as.add_paragraph()
p_as_d.text = (
    "• 한신대 BI 센터 내 입주 공간 확보 및 인프라 지원\n"
    "• H30 디스플레이 시제품 테스트 및 local-first AI 1차 고도화 완료\n"
    "• 지식재산권(특허) 출원 연계 컨설팅 프로세스 지원 진행 중"
)
p_as_d.font.name = "Arial"
p_as_d.font.size = Pt(12)
p_as_d.font.color.rgb = COLOR_WHITE

# TO-BE Card
add_card_bg(slide, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.3), COLOR_GOLD)
tobe_box = slide.shapes.add_textbox(Inches(7.0), Inches(4.45), Inches(5.3), Inches(2.0))
tf_to = tobe_box.text_frame
tf_to.word_wrap = True
tf_to.margin_left = tf_to.margin_top = tf_to.margin_right = tf_to.margin_bottom = 0

p_to_t = tf_to.paragraphs[0]
p_to_t.text = "TO-BE: 향후 확장 비전 (Future Synergy)"
p_to_t.font.name = "Arial"
p_to_t.font.size = Pt(14)
p_to_t.font.bold = True
p_to_t.font.color.rgb = COLOR_GOLD
p_to_t.space_after = Pt(6)

p_to_d = tf_to.add_paragraph()
p_to_d.text = (
    "• 산학 R&D 공동 연구 및 AI/광학 관련 특허 공동 출원 진행\n"
    "• AI 미디어 연산 특화 정부 지원 과제 공동 수주\n"
    "• 플랫폼 인프라 기반의 메타버스/공간 컴퓨팅 고부가가치 응용 확장"
)
p_to_d.font.name = "Arial"
p_to_d.font.size = Pt(12)
p_to_d.font.color.rgb = COLOR_WHITE

# Save the PowerPoint presentation
out_pptx_path = "MBOX_Presentation_WarmDark.pptx"
prs.save(out_pptx_path)
print(f"Presentation successfully compiled to: {out_pptx_path}")
