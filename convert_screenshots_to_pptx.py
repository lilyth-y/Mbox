import os
from pptx import Presentation
from pptx.util import Inches

def main():
    prs = Presentation()
    # Configure 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # We will use the blank slide layout
    blank_layout = prs.slide_layouts[6]
    
    # Path where screenshots are stored
    screenshot_dir = "screenshots"
    if not os.path.exists(screenshot_dir):
        print(f"Error: Directory {screenshot_dir} does not exist. Run capture_slides.py first.")
        return
        
    print("Compiling pixel-perfect PowerPoint from slide screenshots...")
    
    # Sort slide images to ensure sequential ordering (slide_01, slide_02, etc.)
    image_files = sorted([f for f in os.listdir(screenshot_dir) if f.startswith("slide_") and f.endswith(".png")])
    print(f"Found {len(image_files)} screenshots to import.")
    
    for img_name in image_files:
        img_path = os.path.join(screenshot_dir, img_name)
        print(f"Adding {img_name} to PowerPoint...")
        
        # Add blank slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Insert screenshot to span the entire slide (full bleed)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        
    out_pptx = "MBOX_Presentation_WarmDark_PixelPerfect.pptx"
    prs.save(out_pptx)
    print(f"\n[SUCCESS] Pixel-perfect presentation saved successfully to: {out_pptx}")

if __name__ == "__main__":
    main()
